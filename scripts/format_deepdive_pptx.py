"""
format_deepdive_pptx.py — convert a deep-dive .md report into a PowerPoint deck.

Slide map (9 slides):
  1. Cover
  2. Executive Scorecard summary
  3. Business Profile + Capital Efficiency
  4. Forensic Findings
  5. Fraud & Governance Detector
  6. Risk Scorecard (weighted matrix)
  7. Investment Thesis + Earnings Quality
  8. PM One-Pager (verdict)
  9. Data Reliability + Sources

Usage:
    python scripts/format_deepdive_pptx.py <report.md> [--out <output.pptx>]

Or import:
    from format_deepdive_pptx import md_to_pptx
    pptx_bytes = md_to_pptx(md_text, company_name, symbol, isin)
"""
from __future__ import annotations
import os, io, re, sys, argparse, datetime as dt
from typing import Optional

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    sys.exit("python-pptx not installed. Run: pip install python-pptx")

# --------------------------------------------------------------------------
# Colour palette
# --------------------------------------------------------------------------
DARK_BLUE  = RGBColor(0x1A, 0x1A, 0x2E)
MID_BLUE   = RGBColor(0x16, 0x21, 0x3E)
ACCENT     = RGBColor(0x0F, 0x3A, 0x6B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF0, 0xF0, 0xF0)
RED        = RGBColor(0xC0, 0x39, 0x2B)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
AMBER      = RGBColor(0xE6, 0x7E, 0x22)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------
def _set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _textbox(slide, text: str, left, top, width, height,
             font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def _bullet_box(slide, lines: list[str], left, top, width, height,
                font_size=14, color=WHITE, title: str = ""):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    if title:
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.size = Pt(font_size + 2); r.font.bold = True; r.font.color.rgb = color
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if (title or i > 0) else tf.paragraphs[0]
        p.level = 1 if line.startswith("  ") else 0
        r = p.add_run()
        r.text = line.strip().lstrip("•-* ")
        r.font.size = Pt(font_size); r.font.color.rgb = color


# --------------------------------------------------------------------------
# Section extraction
# --------------------------------------------------------------------------
SECTION_MAP = {
    "source_coverage":  [r"A\.\s*Source Coverage"],
    "business_profile": [r"B\.\s*Business Profile"],
    "capital_eff":      [r"C\.\s*Capital Efficiency"],
    "forensic":         [r"D\.\s*Forensic Findings"],
    "fraud_gov":        [r"E\.\s*Fraud.*?Governance"],
    "risk":             [r"F\.\s*Risk Scorecard"],
    "thesis":           [r"G\.\s*Investment Thesis"],
    "earnings_quality": [r"H\.\s*Earnings Quality"],
    "data_reliability": [r"I\.\s*Data Reliability"],
    "analytical_notes": [r"LAYER\s+2.*?ANALYTICAL NOTES", r"ANALYTICAL NOTES"],
    "pm_onepager":      [r"PM ONE[-\s]PAGER"],
}


def _extract_sections(md_text: str) -> dict[str, str]:
    lines = md_text.splitlines()
    sections: dict[str, list[str]] = {k: [] for k in SECTION_MAP}
    current = None
    for line in lines:
        clean = re.sub(r"[=\-#*|]+", " ", line).strip()
        matched = False
        for key, patterns in SECTION_MAP.items():
            for pat in patterns:
                if re.search(pat, clean, re.IGNORECASE):
                    current = key; matched = True; break
            if matched: break
        if not matched and current:
            sections[current].append(line)
    return {k: "\n".join(v).strip() for k, v in sections.items()}


def _bullets(text: str, max_lines: int = 8) -> list[str]:
    """Extract meaningful bullet points from a text block."""
    lines = []
    for l in text.splitlines():
        s = l.strip()
        if not s or re.match(r"^[=\-]{3,}", s):
            continue
        s = re.sub(r"^\*{1,2}|^\#{1,3}\s*", "", s).strip()
        if len(s) > 10:
            lines.append(s)
        if len(lines) >= max_lines:
            break
    return lines


def _verdict_color(text: str) -> RGBColor:
    t = text.upper()
    if "APPROVE" in t or "BUY" in t:    return GREEN
    if "REJECT" in t or "AVOID" in t:   return RED
    return AMBER


# --------------------------------------------------------------------------
# Slide builders
# --------------------------------------------------------------------------
def _slide_cover(prs: Presentation, name: str, symbol: str, isin: str, coverage: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    _set_bg(slide, DARK_BLUE)
    _textbox(slide, "DEEP DIVE REPORT", Inches(0.5), Inches(1.2),
             Inches(12), Inches(1), font_size=32, bold=True, align=PP_ALIGN.CENTER)
    _textbox(slide, name, Inches(0.5), Inches(2.5),
             Inches(12), Inches(1), font_size=28, bold=True, align=PP_ALIGN.CENTER)
    _textbox(slide, f"{symbol}  ·  {isin}", Inches(0.5), Inches(3.4),
             Inches(12), Inches(0.6), font_size=18, align=PP_ALIGN.CENTER, color=LIGHT_GREY)
    _textbox(slide, f"Generated {dt.datetime.now():%Y-%m-%d}",
             Inches(0.5), Inches(4.2), Inches(12), Inches(0.5),
             font_size=14, align=PP_ALIGN.CENTER, color=LIGHT_GREY, italic=True)
    if coverage:
        _textbox(slide, f"Coverage: {coverage}", Inches(0.5), Inches(4.9),
                 Inches(12), Inches(0.5), font_size=13, align=PP_ALIGN.CENTER,
                 color=LIGHT_GREY, italic=True)


def _slide_two_col(prs, title: str, left_title: str, left_text: str,
                   right_title: str, right_text: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MID_BLUE)
    _textbox(slide, title, Inches(0.4), Inches(0.2), Inches(12), Inches(0.6),
             font_size=22, bold=True)
    _bullet_box(slide, _bullets(left_text),  Inches(0.4), Inches(1.0),
                Inches(5.9), Inches(5.8), title=left_title)
    _bullet_box(slide, _bullets(right_text), Inches(6.6), Inches(1.0),
                Inches(6.3), Inches(5.8), title=right_title)


def _slide_single(prs, title: str, section_text: str, accent_line: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MID_BLUE)
    _textbox(slide, title, Inches(0.4), Inches(0.2), Inches(12), Inches(0.6),
             font_size=22, bold=True)
    if accent_line:
        col = _verdict_color(accent_line)
        _textbox(slide, accent_line, Inches(0.4), Inches(0.9), Inches(12), Inches(0.5),
                 font_size=16, bold=True, color=col)
        top_offset = Inches(1.5)
    else:
        top_offset = Inches(1.0)
    _bullet_box(slide, _bullets(section_text, max_lines=12),
                Inches(0.4), top_offset, Inches(12.3), Inches(7.5) - top_offset - Inches(0.3))


def _slide_pm_onepager(prs, text: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_BLUE)
    _textbox(slide, "PM ONE-PAGER", Inches(0.4), Inches(0.2), Inches(12), Inches(0.6),
             font_size=22, bold=True)
    # extract verdict line
    verdict = ""
    for l in text.splitlines():
        if re.search(r"APPROVE|REJECT|WATCHLIST|BUY|AVOID", l, re.IGNORECASE):
            verdict = l.strip(); break
    if verdict:
        col = _verdict_color(verdict)
        _textbox(slide, verdict, Inches(0.4), Inches(1.0), Inches(12), Inches(0.7),
                 font_size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    _bullet_box(slide, _bullets(text, max_lines=14),
                Inches(0.4), Inches(1.9), Inches(12.3), Inches(5.3))


# --------------------------------------------------------------------------
# Main converter
# --------------------------------------------------------------------------
def md_to_pptx(md_text: str, name: str = "", symbol: str = "",
               isin: str = "", coverage: str = "") -> bytes:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    sec = _extract_sections(md_text)

    # Slide 1 — Cover
    _slide_cover(prs, name or "Company", symbol or "", isin or "", coverage)

    # Slide 2 — Executive Scorecard (source coverage summary)
    _slide_single(prs, "Executive Scorecard — Source Coverage", sec["source_coverage"])

    # Slide 3 — Business Profile + Capital Efficiency
    _slide_two_col(prs, "Business Profile & Capital Efficiency",
                   "Business Profile",   sec["business_profile"],
                   "Capital Efficiency", sec["capital_eff"])

    # Slide 4 — Forensic Findings
    _slide_single(prs, "Forensic Findings", sec["forensic"])

    # Slide 5 — Fraud & Governance
    _slide_single(prs, "Fraud & Governance Detector", sec["fraud_gov"])

    # Slide 6 — Risk Scorecard
    _slide_single(prs, "Risk Scorecard", sec["risk"])

    # Slide 7 — Investment Thesis + Earnings Quality
    _slide_two_col(prs, "Investment Thesis & Earnings Quality",
                   "Investment Thesis",  sec["thesis"],
                   "Earnings Quality",   sec["earnings_quality"])

    # Slide 8 — PM One-Pager
    _slide_pm_onepager(prs, sec["pm_onepager"])

    # Slide 9 — Data Reliability + Sources
    _slide_single(prs, "Data Reliability & Sources", sec["data_reliability"])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("input", help="path to .md report file")
    ap.add_argument("--out", help="output .pptx path")
    args = ap.parse_args()

    with open(args.input, encoding="utf-8") as f:
        md_text = f.read()

    m = re.search(r"#\s+Deep Dive.*?([A-Z][^(]+)\s+\((\w+)\s*/\s*(INE\w+)\)", md_text)
    name   = m.group(1).strip() if m else os.path.splitext(os.path.basename(args.input))[0]
    symbol = m.group(2) if m else ""
    isin   = m.group(3) if m else ""

    pptx_bytes = md_to_pptx(md_text, name, symbol, isin)
    out_path = args.out or os.path.splitext(args.input)[0] + ".pptx"
    with open(out_path, "wb") as f:
        f.write(pptx_bytes)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
